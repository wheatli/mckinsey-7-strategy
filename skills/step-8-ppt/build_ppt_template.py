#!/usr/bin/env python3
"""Build strategy report PPT from step-7 recommendation.

Spec (default):
- Font: 微软雅黑 (Microsoft YaHei) for both Latin and East Asian
- Background: white
- Primary color: #1E53A4 (blue)
- Auxiliary: neutral gray (deep #404040 / mid #7F7F7F / light #E6E6E6)
- Highlight: #D80C18 (red) — for key numbers / risk / STOP gates
- Aspect: 16:9 widescreen (13.333" × 7.5")

Usage:
    pip3 install --user --break-system-packages python-pptx
    cp <skill-dir>/build_ppt_template.py docs/strategy/<code>/build_ppt.py
    # 1) edit PROJECT_* constants below
    # 2) fill each slide_* function with content from step-7-recommendation.md
    # 3) python3 build_ppt.py

Cardinal rules (do not violate):
1. ALL text goes through `set_run()` so the East Asian font is bound.
   Do NOT do `text_frame.text = "..."` directly — Chinese will fall back to 宋体.
2. Inside Python string literals, never use paired Chinese double quotes "…".
   They confuse Python's parser. Use 「…」 or 《…》 instead.
3. Do not introduce numbers / conclusions that are not in step-7.
   PPT is a projection, not a re-creation.
4. After every content block, run:
       python3 -c "import ast; ast.parse(open('build_ppt.py').read())"
   to catch syntax errors early.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ============================================================
# PROJECT METADATA — edit these for each project
# ============================================================
PROJECT_CODE = "<project-code>"          # e.g. "futu-web3-strategy"
PROJECT_TITLE = "<战略主标题>"             # e.g. "富途 WEB3 战略建议"
PROJECT_SUBTITLE = "<副标题:主结论一句话>"   # e.g. "HK 旗舰 / SG 延伸 / US 火种"
PROJECT_DECISION_OWNER = "<决策者岗位>"
PROJECT_DATE = "YYYY-MM-DD"
PROJECT_FOOTER = f"{PROJECT_TITLE} · {PROJECT_DATE} · 机密"
OUTPUT_PATH = f"{PROJECT_CODE}-report.pptx"

# ============================================================
# THEME — colors & fonts. Override here if user customizes spec.
# ============================================================
PRIMARY = RGBColor(0x1E, 0x53, 0xA4)
HIGHLIGHT = RGBColor(0xD8, 0x0C, 0x18)
GRAY_DARK = RGBColor(0x40, 0x40, 0x40)
GRAY_MID = RGBColor(0x7F, 0x7F, 0x7F)
GRAY_LIGHT = RGBColor(0xE6, 0xE6, 0xE6)
PRIMARY_BG = RGBColor(0xEA, 0xF1, 0xF9)
HIGHLIGHT_BG = RGBColor(0xFF, 0xF2, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)

FONT_LATIN = "Microsoft YaHei"
FONT_EAST = "Microsoft YaHei"

# ============================================================
# CANVAS
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


# ============================================================
# HELPERS — do not modify unless you know why
# ============================================================
def set_run(run, text, *, size=14, bold=False, color=BLACK, font=FONT_LATIN):
    """Set run text + font face + size + bold + color.

    Critically: also bind the East Asian font face via <a:ea>, so that
    Chinese characters render in 微软雅黑 instead of falling back to 宋体.
    All text in this module MUST go through this function.
    """
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for ea in rPr.findall(qn("a:ea")):
        rPr.remove(ea)
    ea_el = rPr.makeelement(qn("a:ea"), {"typeface": FONT_EAST})
    rPr.append(ea_el)
    for la in rPr.findall(qn("a:latin")):
        rPr.remove(la)
    la_el = rPr.makeelement(qn("a:latin"), {"typeface": FONT_LATIN})
    rPr.append(la_el)


def add_textbox(slide, x, y, w, h, text, *, size=14, bold=False,
                color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return tb


def add_multi_textbox(slide, x, y, w, h, lines, *, default_size=14,
                      anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """lines: list of dict {text, size, bold, color, align, bullet, space_before}."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        p.line_spacing = line_spacing
        if line.get("space_before"):
            p.space_before = Pt(line["space_before"])
        text = line["text"]
        if line.get("bullet"):
            text = "•  " + text
        set_run(p.add_run(), text,
                size=line.get("size", default_size),
                bold=line.get("bold", False),
                color=line.get("color", BLACK))
    return tb


def add_rect(slide, x, y, w, h, *, fill=PRIMARY, line=None, line_w=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    if shape.has_text_frame:
        shape.text_frame.text = ""
    return shape


def add_line(slide, x1, y1, x2, y2, *, color=PRIMARY, weight=2):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def page_footer(slide, page_no, total):
    add_rect(slide, Emu(0), SH - Inches(0.35), SW, Inches(0.35), fill=GRAY_LIGHT)
    add_textbox(slide, Inches(0.4), SH - Inches(0.32), Inches(8), Inches(0.28),
                PROJECT_FOOTER, size=9, color=GRAY_MID)
    add_textbox(slide, SW - Inches(1.5), SH - Inches(0.32), Inches(1.1), Inches(0.28),
                f"{page_no} / {total}", size=9, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def page_header(slide, title, subtitle=None):
    add_rect(slide, Emu(0), Emu(0), Inches(0.25), SH, fill=PRIMARY)
    add_textbox(slide, Inches(0.55), Inches(0.30), Inches(11), Inches(0.55),
                title, size=24, bold=True, color=PRIMARY)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.85), Inches(11), Inches(0.4),
                    subtitle, size=12, color=GRAY_MID)
    add_rect(slide, Inches(0.55), Inches(1.30), Inches(1.2), Inches(0.04),
             fill=HIGHLIGHT)


# ============================================================
# SLIDE 1 — COVER
# Source: step-7 §1 main message (Executive Summary P 部分)
# ============================================================
def slide_cover():
    s = blank_slide()
    add_rect(s, Emu(0), Emu(0), SW, SH, fill=WHITE)
    add_rect(s, Emu(0), Emu(0), Inches(4.8), SH, fill=PRIMARY)
    add_rect(s, Emu(0), Inches(2.8), Inches(4.8), Inches(0.10), fill=HIGHLIGHT)
    # left panel labels
    add_textbox(s, Inches(0.6), Inches(0.6), Inches(4), Inches(0.35),
                "STRATEGY  ·  REPORT", size=12, color=WHITE, bold=True)
    add_textbox(s, Inches(0.6), Inches(1.2), Inches(4), Inches(0.35),
                "麦肯锡七步成诗 · 战略建议书", size=11, color=WHITE)
    add_textbox(s, Inches(0.6), Inches(3.0), Inches(4), Inches(0.4),
                "PROJECT", size=11, color=WHITE, bold=True)
    add_textbox(s, Inches(0.6), Inches(3.45), Inches(4), Inches(0.4),
                PROJECT_CODE, size=14, color=WHITE)
    add_textbox(s, Inches(0.6), Inches(4.1), Inches(4), Inches(0.4),
                "DECISION OWNER", size=11, color=WHITE, bold=True)
    add_textbox(s, Inches(0.6), Inches(4.55), Inches(4), Inches(0.4),
                PROJECT_DECISION_OWNER, size=14, color=WHITE)
    add_textbox(s, Inches(0.6), Inches(5.2), Inches(4), Inches(0.4),
                "DATE", size=11, color=WHITE, bold=True)
    add_textbox(s, Inches(0.6), Inches(5.65), Inches(4), Inches(0.4),
                PROJECT_DATE, size=14, color=WHITE)
    # right main title
    add_textbox(s, Inches(5.3), Inches(2.0), Inches(7.6), Inches(0.6),
                PROJECT_TITLE, size=36, bold=True, color=PRIMARY)
    add_textbox(s, Inches(5.3), Inches(2.85), Inches(7.6), Inches(0.5),
                PROJECT_SUBTITLE, size=20, color=GRAY_DARK)
    # TODO: red highlight box — paste step-7 §1 P 部分(主结论简版,3-4 行)


# ============================================================
# SLIDE 2 — AGENDA
# Source: step-7 章节列表
# ============================================================
def slide_agenda():
    s = blank_slide()
    page_header(s, "议程", "14 页 · 约 25 分钟走查")
    # TODO: 4×4 grid of agenda items, each with number + title


# ============================================================
# SLIDE 3 — EXECUTIVE SUMMARY (SCP)
# Source: step-7 §1 SCP 结构
# ============================================================
def slide_exec_summary():
    s = blank_slide()
    page_header(s, "执行摘要 · SCP 结构", "情境 → 冲突 → 建议")
    top = Inches(1.55)
    col_h = Inches(5.4)
    x0 = Inches(0.55)
    col_w = Inches(4.05)
    gap = Inches(0.15)
    # S column
    add_rect(s, x0, top, col_w, Inches(0.55), fill=GRAY_DARK)
    add_textbox(s, x0 + Inches(0.2), top + Inches(0.10), col_w, Inches(0.4),
                "S · 情境", size=14, bold=True, color=WHITE)
    add_rect(s, x0, top + Inches(0.55), col_w, col_h - Inches(0.55), fill=GRAY_LIGHT)
    # TODO: add_multi_textbox 填 step-7 §1 S
    # C column
    x = x0 + col_w + gap
    add_rect(s, x, top, col_w, Inches(0.55), fill=PRIMARY)
    add_textbox(s, x + Inches(0.2), top + Inches(0.10), col_w, Inches(0.4),
                "C · 冲突", size=14, bold=True, color=WHITE)
    add_rect(s, x, top + Inches(0.55), col_w, col_h - Inches(0.55), fill=GRAY_LIGHT)
    # TODO: add_multi_textbox 填 step-7 §1 C
    # P column
    x = x0 + 2 * (col_w + gap)
    add_rect(s, x, top, col_w, Inches(0.55), fill=HIGHLIGHT)
    add_textbox(s, x + Inches(0.2), top + Inches(0.10), col_w, Inches(0.4),
                "P · 建议", size=14, bold=True, color=WHITE)
    add_rect(s, x, top + Inches(0.55), col_w, col_h - Inches(0.55), fill=HIGHLIGHT_BG)
    # TODO: add_multi_textbox 填 step-7 §1 P (动词起首)


# ============================================================
# SLIDE 4 — THREE REASONS + THREE NUMBERS
# Source: step-7 §1 三大理由 + 关键数字
# ============================================================
def slide_reasons_and_numbers():
    s = blank_slide()
    page_header(s, "三大理由 + 三个关键数字", "建议成立的实证支撑")
    # TODO: 上半 3 栏理由,下半 3 个大字数字(size=32 bold HIGHLIGHT)


# ============================================================
# SLIDE 5 — MARKET LAYERING
# Source: step-7 §2.1 市场分层行 + step-6 论点 A
# ============================================================
def slide_market():
    s = blank_slide()
    page_header(s, "市场分层 · HK / SG / US",
                "18 个月内只有 HK 能零售全栈,资源不能均摊")
    # TODO: 三栏对比 — HK 旗舰 / SG 延伸 / US 火种;每栏列产品集 + 监管限制 + 18 月里程碑


# ============================================================
# SLIDE 6 — PRODUCT ROADMAP
# Source: step-7 §2.1 V1 / V1.5 行 + 论点 B
# ============================================================
def slide_product():
    s = blank_slide()
    page_header(s, "产品节奏 · V1 → V1.5 + 整合护城河",
                "分阶段交付 + RWA / 钱包差异化 + 不自研 L2")
    # TODO: 时间线 0 → 6m → 12m,标 V1 / V1.5 / RWA / 钱包 anchor


# ============================================================
# SLIDE 7 — GROWTH ENGINE RELAY
# Source: step-6 论点 C + step-7 §2.1 增长引擎行
# ============================================================
def slide_growth():
    s = blank_slide()
    page_header(s, "增长引擎接力 · 存量 → 外部",
                "前 6-9 月 cross-sell 主引擎,9 月节点必接外部")
    # TODO: 左半 CAC 对比柱 (存量 $10-20 vs 外部 $30-100);右半 6-9 月接力时间线


# ============================================================
# SLIDE 8 — WHAT WE DO NOT RECOMMEND
# Source: step-7 §2.2 排除清单
# ============================================================
def slide_dont_do():
    s = blank_slide()
    page_header(s, "我们不建议什么", "提前回答决策者会问的反问题")
    # TODO: 7 行表格 (备选 / 排除理由 / 来源)
    # 表头:add_rect 主色填充 + add_textbox 白字
    # 数据行:交替 GRAY_LIGHT / WHITE 底色


# ============================================================
# SLIDE 9 — IMPLEMENTATION STAGES (with GO/STOP gates)
# Source: step-7 §3 全表
# ============================================================
def slide_stages():
    s = blank_slide()
    page_header(s, "实施路径 · 4 阶段决策门",
                "每阶段 GO/STOP 条件明确,不留单程票")
    # TODO: 4 阶段并排 (准备 / V1+试点 / V1.5+护城河 / 复盘+L2 重审)
    # 每个阶段:阶段名(主色块) / 时间 / 关键动作 / GO 条件(绿) / STOP 条件(红 HIGHLIGHT)


# ============================================================
# SLIDE 10 — RESOURCES (phased)
# Source: step-7 §4 全表
# ============================================================
def slide_resources():
    s = blank_slide()
    page_header(s, "资源需求 · 分阶段投入",
                "决策门未过则停止下一阶段投入,不一次 all-in")
    # TODO: 双行表 — 上行资金,下行 FTE;按 4 阶段列展开


# ============================================================
# SLIDE 11 — RISKS WITH EARLY SIGNALS
# Source: step-7 §5 全表
# ============================================================
def slide_risk():
    s = blank_slide()
    page_header(s, "风险与早期信号 · 7 项",
                "每项风险都有可观测的早期触发信号")
    # TODO: 7 行风险表 (R# / 风险 / 概率 / 影响 / 早期信号)
    # 概率/影响列用色块标识等级;早期信号列阈值用 HIGHLIGHT


# ============================================================
# SLIDE 12 — KPI MATRIX
# Source: step-7 §6 全表
# ============================================================
def slide_kpi():
    s = blank_slide()
    page_header(s, "KPI 矩阵 · 业务 / 财务 / 客户 / 团队",
                "≥ 6 个先行指标,30d / 90d / 12m / 18m 目标")
    # TODO: 4 行(维度) × 5 列(指标 / 当前 / 30d / 90d / 12m / 18m)


# ============================================================
# SLIDE 13 — NEXT 7 DAYS
# Source: step-7 §7 全表
# ============================================================
def slide_next_7_days():
    s = blank_slide()
    page_header(s, "Next 7 Days · 立项启动后第一周",
                "假设决策者今天批准,第一周必须发生的具体动作")
    # TODO: 7 张日卡 (Day 1 ~ Day 7),每卡列具体会议 / 邮件 / 文件 / 决定


# ============================================================
# SLIDE 14 — DECISION REQUEST
# Source: step-7 §1 P + §3 决策门 综合
# ============================================================
def slide_decision():
    s = blank_slide()
    page_header(s, "决策请求", "本次会议需要 BU 总经理给出 3 个 GO 决定")
    # TODO: 3 个 GO 决策 + 立项启动会议时点(Day 1)
    # 用 HIGHLIGHT 色块标 GO 字眼


# ============================================================
# MAIN
# ============================================================
def main():
    slide_funcs = [
        slide_cover, slide_agenda, slide_exec_summary, slide_reasons_and_numbers,
        slide_market, slide_product, slide_growth, slide_dont_do,
        slide_stages, slide_resources, slide_risk, slide_kpi,
        slide_next_7_days, slide_decision,
    ]
    total = len(slide_funcs)
    for i, fn in enumerate(slide_funcs, 1):
        fn()
        if i > 1:  # cover has no footer
            page_footer(prs.slides[i - 1], i, total)

    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")
    print(f"Slides: {total}")


if __name__ == "__main__":
    main()
